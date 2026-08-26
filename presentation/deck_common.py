"""Shared typography and layout helpers for the meeting decks."""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

OUT = None

FONT = "Times New Roman"

INK    = RGBColor(0x1A, 0x1A, 0x1A)
MUTED  = RGBColor(0x5F, 0x5F, 0x5F)
ACCENT = RGBColor(0x1F, 0x4E, 0x79)
GOOD   = RGBColor(0x1E, 0x6B, 0x3C)
WARN   = RGBColor(0x9B, 0x35, 0x25)
AMBER  = RGBColor(0xA3, 0x6E, 0x0A)
HDRBG  = RGBColor(0x1F, 0x4E, 0x79)
ROWBG  = RGBColor(0xF2, 0xF4, 0xF7)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
PALE   = RGBColor(0xEC, 0xEF, 0xF4)
RULE   = RGBColor(0xB8, 0xC2, 0xCE)

# decomposition figure palette: two blues mean "same problem", brick means "different problem"
C_DONE  = RGBColor(0x1F, 0x4E, 0x79)
C_REACH = RGBColor(0xA8, 0xC0, 0xD8)
C_FEAT  = RGBColor(0x9B, 0x45, 0x35)

prs = Presentation()
prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
BLANK = prs.slide_layouts[6]


def _style(par, size, color, bold):
    """Apply typography to a paragraph and every run inside it."""
    par.font.size = Pt(size); par.font.color.rgb = color
    par.font.bold = bold; par.font.name = FONT
    for r in par.runs:
        r.font.size = Pt(size); r.font.color.rgb = color
        r.font.bold = bold; r.font.name = FONT


def slide():
    return prs.slides.add_slide(BLANK)


def title(s, text, sub=None):
    b = s.shapes.add_textbox(Inches(0.6), Inches(0.34), Inches(12.2), Inches(0.9))
    tf = b.text_frame; tf.word_wrap = True
    _style(tf.paragraphs[0], 26, ACCENT, True)
    tf.paragraphs[0].text = text
    _style(tf.paragraphs[0], 26, ACCENT, True)
    if sub:
        q = tf.add_paragraph(); q.text = sub
        _style(q, 12.5, MUTED, False)
    # thin rule under the header
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.62), Inches(1.26), Inches(12.1), Inches(0.018))
    ln.fill.solid(); ln.fill.fore_color.rgb = RULE
    ln.line.fill.background(); ln.shadow.inherit = False
    return b


def bullets(s, items, top, size=13, left=0.65, width=12.1, gap=8, bottom=7.1):
    b = s.shapes.add_textbox(Inches(left), Inches(top), Inches(width),
                             Inches(max(0.4, min(5.4, bottom - top))))
    tf = b.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = it
        _style(p, size, INK, False)
        p.space_after = Pt(gap)
    return b


def table(s, rows, top, left=0.65, width=12.1, height=None, col_w=None,
          font=11.5, highlight=(), hcolor=None):
    nr, nc = len(rows), len(rows[0])
    shp = s.shapes.add_table(nr, nc, Inches(left), Inches(top),
                             Inches(width), Inches(height or 0.32 * nr))
    tbl = shp.table
    tbl.first_row = True
    if col_w:
        tot = sum(col_w)
        for j, w in enumerate(col_w):
            tbl.columns[j].width = Inches(width * w / tot)
    for i, row in enumerate(rows):
        for j, v in enumerate(row):
            c = tbl.cell(i, j); c.text = str(v)
            c.margin_left = c.margin_right = Inches(0.08)
            c.margin_top = c.margin_bottom = Inches(0.03)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            hdr = (i == 0)
            c.fill.solid()
            c.fill.fore_color.rgb = (HDRBG if hdr else
                                     (hcolor or PALE) if i in highlight else
                                     (WHITE if i % 2 else ROWBG))
            bold = hdr or (i in highlight)
            for pa in c.text_frame.paragraphs:
                pa.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
                _style(pa, font, WHITE if hdr else INK, bold)
    return shp


def cell_color(shp, row, col, color, bold=True):
    c = shp.table.cell(row, col)
    for pa in c.text_frame.paragraphs:
        pa.font.color.rgb = color; pa.font.bold = bold; pa.font.name = FONT
        for r in pa.runs:
            r.font.color.rgb = color; r.font.bold = bold; r.font.name = FONT


def link_cell(shp, row, col, url):
    c = shp.table.cell(row, col)
    for pa in c.text_frame.paragraphs:
        for r in pa.runs:
            r.hyperlink.address = url


def box(s, x, y, w, h, text, fill, fg=WHITE, size=12, bold=True,
        shape=MSO_SHAPE.RECTANGLE, line=None):
    sh = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line; sh.line.width = Pt(0.75)
    sh.shadow.inherit = False
    tf = sh.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.text = text; p.alignment = PP_ALIGN.CENTER
    _style(p, size, fg, bold)
    return sh


def label(s, x, y, w, text, size=11, color=INK, bold=False, align=PP_ALIGN.CENTER, h=0.34):
    b = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = b.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.alignment = align
    _style(p, size, color, bold)
    return b


def takeaway(s, text, color=ACCENT, top=6.55):
    b = s.shapes.add_textbox(Inches(0.65), Inches(top), Inches(12.1), Inches(0.7))
    tf = b.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    _style(p, 13.5, color, True)
    return b


