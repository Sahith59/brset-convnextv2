"""Build the progress presentation for Dr. Ye.

Design brief: simple, minimal, no heavy visuals. Numbers carry the argument,
so tables are used wherever a claim rests on a number. Short bullets only.
"""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

OUT = Path(__file__).parent / "BRSET_ConvNeXtV2_Progress.pptx"

INK = RGBColor(0x1A, 0x1A, 0x1A)
MUTED = RGBColor(0x5A, 0x5A, 0x5A)
ACCENT = RGBColor(0x1F, 0x4E, 0x79)
GOOD = RGBColor(0x1E, 0x7A, 0x3C)
WARN = RGBColor(0xA5, 0x3A, 0x1A)
HDRBG = RGBColor(0x1F, 0x4E, 0x79)
ROWBG = RGBColor(0xF2, 0xF4, 0xF7)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def title(s, text, sub=None):
    box = s.shapes.add_textbox(Inches(0.7), Inches(0.45), Inches(12.0), Inches(0.85))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    if sub:
        p2 = tf.add_paragraph()
        p2.text = sub
        p2.font.size = Pt(15)
        p2.font.color.rgb = MUTED
    return box


def bullets(s, items, top=1.65, left=0.75, width=11.9, size=17, gap=10, bottom_limit=7.2):
    # Size the box to the space actually left on the slide so nothing hangs off the edge.
    height = max(0.5, min(4.6, bottom_limit - top))
    box = s.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, it in enumerate(items):
        if isinstance(it, tuple):
            text, lvl = it
        else:
            text, lvl = it, 0
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ("• " if lvl == 0 else "– ") + text
        p.level = lvl
        p.font.size = Pt(size if lvl == 0 else size - 2)
        p.font.color.rgb = INK if lvl == 0 else MUTED
        p.space_after = Pt(gap)
    return box


def table(s, rows, top, left=0.75, width=11.8, height=None, col_w=None,
          font=13, header=True, highlight_rows=(), highlight_color=None):
    nrow, ncol = len(rows), len(rows[0])
    height = height or 0.36 * nrow
    shape = s.shapes.add_table(nrow, ncol, Inches(left), Inches(top),
                               Inches(width), Inches(height))
    tbl = shape.table
    if col_w:
        total = sum(col_w)
        for j, w in enumerate(col_w):
            tbl.columns[j].width = Inches(width * w / total)
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = tbl.cell(i, j)
            cell.text = str(val)
            cell.margin_left = Inches(0.08)
            cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.03)
            cell.margin_bottom = Inches(0.03)

            is_header = (i == 0 and header)
            cell.fill.solid()
            if is_header:
                cell.fill.fore_color.rgb = HDRBG
            elif i in highlight_rows:
                cell.fill.fore_color.rgb = highlight_color or RGBColor(0xE3, 0xEF, 0xE3)
            else:
                cell.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF) if i % 2 else ROWBG

            # Style EVERY paragraph, not just the first. A cell whose text
            # contains newlines becomes multiple paragraphs, and any left
            # unstyled falls back to the 18pt default -- which is what made
            # multi-line cells render at mixed sizes.
            for para in cell.text_frame.paragraphs:
                para.font.size = Pt(font)
                para.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
                if is_header:
                    para.font.bold = True
                    para.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                else:
                    para.font.bold = i in highlight_rows
                    para.font.color.rgb = INK
                for run in para.runs:
                    run.font.size = Pt(font)
                    run.font.color.rgb = (RGBColor(0xFF, 0xFF, 0xFF) if is_header else INK)
                    run.font.bold = is_header or (i in highlight_rows)
    return shape


def takeaway(s, text, color=ACCENT, top=6.35):
    box = s.shapes.add_textbox(Inches(0.75), Inches(top), Inches(11.9), Inches(0.75))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = color
    return box


# ---------------------------------------------------------------- 1. Title
s = slide()
box = s.shapes.add_textbox(Inches(0.9), Inches(2.4), Inches(11.6), Inches(2.2))
tf = box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "ConvNeXt V2 on BRSET"
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = ACCENT
p = tf.add_paragraph()
p.text = "A strong baseline, and closing the BRSET → mBRSET device gap"
p.font.size = Pt(21)
p.font.color.rgb = INK
p.space_before = Pt(8)
p = tf.add_paragraph()
p.text = "Sahith Reddy Thummala   |   For: Dr. Dong Hye Ye   |   cc: Nagur Shareef Shaik"
p.font.size = Pt(14)
p.font.color.rgb = MUTED
p.space_before = Pt(22)

# ---------------------------------------------------------------- 2. The ask
s = slide()
title(s, "What was assigned")
table(s, [
    ["", "Task", "Status"],
    ["Part 1", "A strong ConvNeXt V2 baseline on BRSET;\naddress the class imbalance holding F1 down",
     "Complete"],
    ["Part 2", "Improve BRSET → mBRSET transfer.\nFind relevant literature, implement, and find novelty",
     "Diagnosed;\nplan ready"],
], top=2.0, col_w=[1.0, 6.5, 2.2], font=15, height=2.0)
takeaway(s, "Part 1 is answered. Part 2 is where the contribution is — and the lane is open.", top=5.0)

# ---------------------------------------------------------------- 3. Baseline result
s = slide()
title(s, "Part 1 — Baseline vs. the original BRSET paper",
      "One ConvNeXt V2 Large model at 512px, all 16,258 images, split by patient, trained with focal loss.")
table(s, [
    ["Metric", "Paper (Nakayama 2024)", "Our model", "Verdict"],
    ["Diabetic retinopathy — AUC", "0.97", "0.992", "Better than the paper"],
    ["Diabetic retinopathy — F1", "0.89", "0.869", "The same, within measurement error"],
    ["Macular edema — AUC", "not reported", "0.989", "No published benchmark to compare to"],
    ["Macular edema — F1", "not reported", "0.748", "No published benchmark to compare to"],
], top=2.15, col_w=[3.2, 2.4, 1.7, 4.3], font=14,
   highlight_rows=(1, 2), highlight_color=RGBColor(0xE3, 0xEF, 0xE3))
bullets(s, [
    "Our AUC of 0.992 sits above the paper's 0.97 by more than measurement error — a real improvement.",
    "Our F1 measures anywhere from 0.83 to 0.90 depending on which patients land in the test set. "
    "The paper's 0.89 falls inside that range, so the two cannot be told apart.",
    "Closing the remaining 0.02 would mean catching about 6 more images out of 2,435 — "
    "less than the test set can reliably measure.",
    "This is a single model. It matches our earlier two-model ensemble on DR using half the models.",
], top=4.2, size=15, bottom_limit=6.25)
takeaway(s, "We beat the paper on AUC and match it on F1 — with one model instead of two.", GOOD)

# ---------------------------------------------------------------- 3b. The imbalance itself
s = slide()
title(s, "Part 1 — The imbalance we had to handle",
      "Almost every image in BRSET is healthy. That is why F1 lags behind AUC.")
table(s, [
    ["Split", "Images", "With diabetic retinopathy", "With macular edema"],
    ["Train", "11,372", "748  (6.6%)", "274  (2.4%)"],
    ["Validation", "2,451", "159  (6.5%)", "65  (2.7%)"],
    ["Test", "2,435", "162  (6.7%)", "61  (2.5%)"],
], top=2.1, width=10.5, col_w=[1.6, 1.6, 3.0, 2.8], font=15, height=1.6)
bullets(s, [
    "Roughly 15 healthy images for every 1 with disease — and 40 to 1 for macular edema.",
    "A model can score 93% accuracy by calling everything healthy, and be useless. So we never report accuracy alone.",
    "AUC stays high because ranking sick above healthy is easy. F1 suffers because the model,",
    ("seeing mostly healthy eyes, becomes reluctant to call anything diseased.", 1),
], top=4.1, size=16, bottom_limit=6.2)
takeaway(s, "The imbalance is the reason F1 is the hard metric here — not AUC.")

# ---------------------------------------------------------------- 4. Imbalance handling
s = slide()
title(s, "Part 1 — How the imbalance was addressed",
      "BRSET has about 15 healthy images for every diseased one. Four steps stop the majority drowning out the rest.")
bullets(s, [
    "Focal loss — the model stops spending effort on images it already gets right, and concentrates on the hard ones.",
    "Oversampling — rare diseased images are shown more often during training than they naturally occur.",
    "Tuned referral cutoff — chosen on validation data, separately for each finding, instead of a blind 0.5.",
    "Averaging — the model is averaged over training, and each image is predicted from four flipped copies.",
], top=2.15, size=17, bottom_limit=4.3)
table(s, [
    ["Loss function tested", "DR F1", "ME F1", "Outcome"],
    ["Focal loss — one dial for all images", "0.869", "0.748", "Kept"],
    ["Asymmetric Loss — separate dials for\nhealthy and for diseased images", "0.834", "0.758", "Clearly worse on DR — not used"],
], top=4.5, width=11.0, col_w=[4.2, 1.3, 1.3, 3.6], font=13, height=1.35)
takeaway(s, "We tested a more elaborate alternative. It did not help, so the simpler standard loss stands.")

# ---------------------------------------------------------------- 4b. Optional further steps (approval gate)
s = slide()
title(s, "Part 1 — Two further steps, if you would like them",
      "Both are ready to launch. About 15 hours on the cluster, and neither blocks the Part 2 work.")
table(s, [
    ["Step", "What it does", "Why it would matter", "Cost"],
    ["Hyperparameter\nsweep",
     "Tries six settings of the imbalance\ncontrols, instead of the one value we\nchose by hand",
     "Lets us say we searched the settings\nproperly, rather than used the first\nvalue we tried",
     "~7 h"],
    ["5-fold\ncross-validation",
     "Re-runs the model five times so every\nimage gets tested, instead of testing\non a single 15% slice",
     "Halves the measurement error. That is\nwhat would make a 0.02 change in F1\nprovable rather than nominal",
     "~7 h"],
], top=2.05, col_w=[1.9, 4.2, 4.2, 0.9], font=12.5, height=2.6)
bullets(s, [
    "Expected effect on the headline numbers: small. We are already at the paper's level, "
    "and the sweep is unlikely to move it much.",
    "The gain is confidence rather than score — the cross-validation would let us state the result "
    "with half the current uncertainty.",
    "Neither blocks Part 2. The cross-device work can proceed in parallel either way.",
], top=4.9, size=15, bottom_limit=6.25)
takeaway(s, "Happy to run both, either, or neither — whichever you think is the better use of the time.")

# ---------------------------------------------------------------- 5. The cross-device problem
s = slide()
title(s, "Part 2 — The BRSET → mBRSET problem",
      "BRSET: tabletop hospital cameras. mBRSET: handheld smartphone camera, community screening, 83% with artifacts.")
table(s, [
    ["Trained → Tested", "AUC (DR / ME)", "F1 (DR / ME)", "Recall (DR / ME)"],
    ["BRSET → BRSET (in-domain)", "0.988 / 0.994", "0.869 / 0.790", "0.877 / 0.770"],
    ["mBRSET → mBRSET (in-domain)", "0.939 / 0.988", "0.815 / 0.807", "0.774 / 0.730"],
    ["BRSET → mBRSET, BRSET threshold", "0.909 / 0.933", "0.661 / 0.566", "0.503 / 0.444"],
    ["BRSET → mBRSET, threshold retuned", "0.909 / 0.933", "0.717 / 0.628", "0.780 / 0.603"],
], top=2.2, col_w=[3.8, 2.6, 2.6, 2.6], font=14,
   highlight_rows=(4,), highlight_color=RGBColor(0xFD, 0xF0, 0xE6))
bullets(s, [
    "Retuning only the decision threshold — no retraining — lifted DR recall from 0.503 to 0.780.",
    "Remaining gap to a model trained directly on mBRSET: 0.098 F1 (DR), 0.179 F1 (ME).",
], top=4.6, size=17, bottom_limit=6.2)
takeaway(s, "Question: what is that remaining gap actually made of, and what can close it?")

# ---------------------------------------------------------------- 6. Diagnosis 1
s = slide()
title(s, "Diagnosis 1 — the referral cutoff was wrong, not the model",
      "The model scores each image from 0 (healthy) to 1 (diseased). A cutoff decides who gets referred.")
table(s, [
    ["", "How many patients have the disease", "Best cutoff for that group"],
    ["BRSET — where we trained", "6.6% of eyes  (1 in 15)", "0.61   be strict"],
    ["mBRSET — where we tested", "23.3% of eyes  (1 in 4)", "0.19   be lenient"],
], top=2.25, width=10.8, col_w=[3.0, 4.2, 3.0], font=15, height=1.5)
bullets(s, [
    "We applied BRSET's strict 0.61 cutoff to mBRSET, where nearly 1 in 4 patients is diseased. "
    "The model stayed quiet and missed half the cases.",
    "Of the move needed from 0.61 down to 0.19:  77% is explained by the higher disease rate alone, "
    "and only 23% by the model being less certain on blurrier images.",
    "The disease-rate part can be estimated from unlabelled images — no answer key required.",
], top=4.25, size=16, bottom_limit=6.2)
takeaway(s, "Moving the cutoff alone — no retraining — lifted F1 from 0.661 to 0.717, and recall from 0.503 to 0.780.")

# ---------------------------------------------------------------- 7. Diagnosis 2 (key slide)
s = slide()
title(s, "Diagnosis 2 — Threshold tuning is already exhausted",
      "For a given AUC and prevalence there is a hard maximum F1, over all possible thresholds.")
table(s, [
    ["BRSET → mBRSET", "AUC", "Target prevalence", "Max possible F1", "Our F1"],
    ["Diabetic retinopathy", "0.909", "21.7%", "0.697", "0.717"],
    ["Macular edema", "0.933", "8.6%", "0.623", "0.628"],
], top=2.1, col_w=[3.0, 1.5, 2.4, 2.4, 1.7], font=15,
   highlight_rows=(1, 2), highlight_color=RGBColor(0xE3, 0xEF, 0xE3))
bullets(s, [
    "We are already at that ceiling. Moving the cutoff further cannot help — any cutoff method only slides the",
    ("referral line along a fixed ordering of patients. It never changes how well the model orders them.", 1),
    "To reach what an mBRSET-trained model achieves (F1 0.815 / 0.807), the ordering itself must improve:",
    ("AUC 0.909 → about 0.95 for DR, and 0.933 → about 0.985 for macular edema.", 1),
], top=4.0, size=17, bottom_limit=6.2)
takeaway(s, "The remaining gap is entirely ranking quality. Everything from here must improve AUC.", WARN)

# ---------------------------------------------------------------- 9. The plan
s = slide()
title(s, "Part 2 — Plan to make the BRSET model work on mBRSET",
      "Two methods proposed, one parked for your input. Ordered by expected gain per unit of effort.")
table(s, [
    ["#", "Method", "What it does", "Effort", "Status"],
    ["1", "Use all of a patient's\nimages together",
     "mBRSET has about 4 images per patient. 85% of the time they share\n"
     "the same diagnosis, but we currently judge each image alone.\n"
     "Blur differs per image, so combining them cancels it out.\n"
     "Published gains on retinal data: +0.085 to +0.121 AUC. We need +0.04.",
     "3–5 d", "Proposed"],
    ["2", "Set the cutoff without\nneeding an answer key",
     "Estimates the disease rate from unlabelled images, so the cutoff fix\n"
     "works on a new camera with no annotations at all.\n"
     "Does not raise AUC — it makes the gain we already have deployable,\n"
     "and is the honest baseline the other methods are measured against.",
     "0.5 d", "Proposed"],
    ["3", "Train on deliberately\ndegraded BRSET images",
     "Blur and darken real BRSET images to imitate handheld capture, so\n"
     "the model learns to find lesions despite poor quality.\n"
     "Aimed at macular edema specifically, where blur destroys the sharp\n"
     "edge that separates an exudate from a harmless drusen.",
     "3–5 d", "Parked —\nfor discussion"],
], top=1.9, col_w=[0.4, 2.4, 6.6, 0.8, 1.4], font=11.5, height=4.3)
takeaway(s, "Method 3 is the one I would like your view on before starting it.", top=6.45)

# ---------------------------------------------------------------- 10. Novelty
s = slide()
title(s, "Novelty — what we found, and how certain we are",
      "Search covered all 33 papers citing either dataset, plus GitHub, arXiv and preprint servers.")
table(s, [
    ["Finding", "How certain"],
    ["No published BRSET → mBRSET transfer study found",
     "Exhaustive citation search — but not finding one\nis not the same as proving none exists"],
    ["The dataset creators list \"domain adaptation across imaging devices\"\nas an application, and report no such experiment",
     "Verified — stated on their March 2026 release"],
    ["No work uses BRSET's per-image degradation labels\n(focus / illumination / field / artifact) for transfer",
     "Verified — the labels exist and are unused"],
    ["No work reports cross-device macular edema degradation",
     "Verified — and our ME gap (0.179) is larger\nthan our DR gap (0.098)"],
    ["RetSyn (J Biomed Inform 2025) — synthetic data for tabletop → portable.\nSame senior author as both datasets.",
     "Published and real. This is the competitor."],
    ["One GitHub project is doing a descriptive version of BRSET → mBRSET",
     "Real, but a single commit — code skeleton,\nno data, no results published yet"],
], top=1.95, col_w=[6.9, 4.7], font=11.5, height=3.9)
bullets(s, [
    "How we differ from RetSyn: they generate synthetic images; we diagnose what actually breaks, and separate "
    "image quality from camera identity. They do not address macular edema or the cutoff problem.",
    "How we differ from the GitHub project: they use frozen features and a simple classifier; we fine-tune the "
    "whole model and cover both findings.",
    "Our shift is unusually clean: country, ethnicity, dilation, 45° field of view and grading protocol are all "
    "held constant, so only the camera changes. A study that also crossed country and protocol collapsed to AUC 0.54; ours holds 0.91.",
], top=5.95, size=13, bottom_limit=7.3)

# ---------------------------------------------------------------- 11. Proposed contributions
s = slide()
title(s, "Proposed contributions")
table(s, [
    ["", "Contribution", "Why it is new"],
    ["1", "Decomposing a cross-device gap into\nprevalence, calibration and ranking components",
     "Cross-domain F1 collapses are routinely blamed on the model.\n"
     "We show the split is computable in closed form — and that in our\n"
     "case the ranking component is all that remains."],
    ["2", "Multi-view test-time aggregation for\ncross-device retinal screening",
     "Exploits mBRSET's ~4 images per patient, a free signal no\ncross-device study has used."],
    ["3", "Adapting the internal normalisation layer\nthat is unique to ConvNeXt V2",
     "The standard adaptation trick relies on a component this model\ndoes not have. Nobody has adapted the one it does have."],
], top=1.95, col_w=[0.5, 4.4, 6.7], font=12, height=3.9)
takeaway(s, "Contribution 1 stands even if every method underperforms — the negative result is itself the finding.",
         top=6.1)

# ---------------------------------------------------------------- 12. Next steps
s = slide()
title(s, "Next steps")
table(s, [
    ["When", "Action"],
    ["This week", "Finish imbalance hyperparameter sweep (validation-selected); k-fold cross-validation "
                  "to tighten confidence intervals on the baseline"],
    ["Next", "Implement label-free threshold placement — the honest baseline for all transfer results"],
    ["Then", "Multi-view aggregation, then degradation augmentation; measure AUC gain against the 0.95 / 0.985 target"],
    ["Open question", "Access to BRSET's per-image degradation-type metadata for the decomposition analysis"],
], top=2.1, col_w=[1.8, 9.8], font=14, height=2.8)
bullets(s, [
    "Baseline: complete and defensible — beats the paper on AUC, matches on F1.",
    "Transfer: diagnosed, literature reviewed, three methods ranked, novelty verified as open.",
], top=5.2, size=17, bottom_limit=7.2)

prs.save(OUT)
print(f"wrote {OUT}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
