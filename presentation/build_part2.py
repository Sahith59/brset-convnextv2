"""Meeting deck: Part 1 result (closed) and Part 2 diagnosis with proposed direction.

Design brief. A reader who has not seen this work should follow every slide
without narration. Numbers and figures carry the argument. Every figure here was
recomputed from the files in results/ rather than carried over on trust.

Typography is Times New Roman throughout. No decorative glyphs, no em dashes.
First person singular, since this is one person's work.
"""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

OUT = Path(__file__).parent / "BRSET_Part2_LiteratureReview.pptx"

FONT = "Times New Roman"

INK    = RGBColor(0x1A, 0x1A, 0x1A)
MUTED  = RGBColor(0x5F, 0x5F, 0x5F)
ACCENT = RGBColor(0x1F, 0x4E, 0x79)
GOOD   = RGBColor(0x1E, 0x6B, 0x3C)
WARN   = RGBColor(0x9B, 0x35, 0x25)
AMBER  = RGBColor(0xA3, 0x6E, 0x0A)
HDRBG  = RGBColor(0x1F, 0x4E, 0x79)
ROWBG  = RGBColor(0xF2, 0xF4, 0xF7)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
PALE   = RGBColor(0xEC, 0xEF, 0xF4)
RULE   = RGBColor(0xB8, 0xC2, 0xCE)

# decomposition figure palette: two blues mean "same problem", brick means "different problem"
C_DONE  = RGBColor(0x1F, 0x4E, 0x79)
C_REACH = RGBColor(0xA8, 0xC0, 0xD8)
C_FEAT  = RGBColor(0x9B, 0x45, 0x35)

prs = Presentation()
prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
BLANK = prs.slide_layouts[6]


def _style(par, size, color, bold):
    """Apply typography to a paragraph and every run inside it."""
    par.font.size = Pt(size); par.font.color.rgb = color
    par.font.bold = bold; par.font.name = FONT
    for r in par.runs:
        r.font.size = Pt(size); r.font.color.rgb = color
        r.font.bold = bold; r.font.name = FONT


def slide():
    return prs.slides.add_slide(BLANK)


def title(s, text, sub=None):
    b = s.shapes.add_textbox(Inches(0.6), Inches(0.34), Inches(12.2), Inches(0.9))
    tf = b.text_frame; tf.word_wrap = True
    _style(tf.paragraphs[0], 26, ACCENT, True)
    tf.paragraphs[0].text = text
    _style(tf.paragraphs[0], 26, ACCENT, True)
    if sub:
        q = tf.add_paragraph(); q.text = sub
        _style(q, 12.5, MUTED, False)
    # thin rule under the header
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.62), Inches(1.26), Inches(12.1), Inches(0.018))
    ln.fill.solid(); ln.fill.fore_color.rgb = RULE
    ln.line.fill.background(); ln.shadow.inherit = False
    return b


def bullets(s, items, top, size=13, left=0.65, width=12.1, gap=8, bottom=7.1):
    b = s.shapes.add_textbox(Inches(left), Inches(top), Inches(width),
                             Inches(max(0.4, min(5.4, bottom - top))))
    tf = b.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = it
        _style(p, size, INK, False)
        p.space_after = Pt(gap)
    return b


def table(s, rows, top, left=0.65, width=12.1, height=None, col_w=None,
          font=11.5, highlight=(), hcolor=None):
    nr, nc = len(rows), len(rows[0])
    shp = s.shapes.add_table(nr, nc, Inches(left), Inches(top),
                             Inches(width), Inches(height or 0.32 * nr))
    tbl = shp.table
    tbl.first_row = True
    if col_w:
        tot = sum(col_w)
        for j, w in enumerate(col_w):
            tbl.columns[j].width = Inches(width * w / tot)
    for i, row in enumerate(rows):
        for j, v in enumerate(row):
            c = tbl.cell(i, j); c.text = str(v)
            c.margin_left = c.margin_right = Inches(0.08)
            c.margin_top = c.margin_bottom = Inches(0.03)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            hdr = (i == 0)
            c.fill.solid()
            c.fill.fore_color.rgb = (HDRBG if hdr else
                                     (hcolor or PALE) if i in highlight else
                                     (WHITE if i % 2 else ROWBG))
            bold = hdr or (i in highlight)
            for pa in c.text_frame.paragraphs:
                pa.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
                _style(pa, font, WHITE if hdr else INK, bold)
    return shp


def cell_color(shp, row, col, color, bold=True):
    c = shp.table.cell(row, col)
    for pa in c.text_frame.paragraphs:
        pa.font.color.rgb = color; pa.font.bold = bold; pa.font.name = FONT
        for r in pa.runs:
            r.font.color.rgb = color; r.font.bold = bold; r.font.name = FONT


def link_cell(shp, row, col, url):
    c = shp.table.cell(row, col)
    for pa in c.text_frame.paragraphs:
        for r in pa.runs:
            r.hyperlink.address = url


def box(s, x, y, w, h, text, fill, fg=WHITE, size=12, bold=True,
        shape=MSO_SHAPE.RECTANGLE, line=None):
    sh = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line; sh.line.width = Pt(0.75)
    sh.shadow.inherit = False
    tf = sh.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.text = text; p.alignment = PP_ALIGN.CENTER
    _style(p, size, fg, bold)
    return sh


def label(s, x, y, w, text, size=11, color=INK, bold=False, align=PP_ALIGN.CENTER, h=0.34):
    b = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = b.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.alignment = align
    _style(p, size, color, bold)
    return b


def takeaway(s, text, color=ACCENT, top=6.55):
    b = s.shapes.add_textbox(Inches(0.65), Inches(top), Inches(12.1), Inches(0.7))
    tf = b.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    _style(p, 13.5, color, True)
    return b


# ═════════════════════════════════════════════ 1. Part 1 final results
s = slide()
title(s, "Part 1: final 40-epoch results",
      "All four runs are identical except for the loss and the sampler. Same backbone, 512 px, "
      "lr 5e-5, 40 epochs, effective batch 64, EMA with 4-way flip TTA, seed 0, "
      "and the same 2,435-image test set.")

table(s, [
    ["Configuration", "DR AUC", "DR F1\n(class-macro)", "DR cases\nmissed", "ME AUC", "ME F1\n(class-macro)"],
    ["Focal with oversampling (baseline in earlier deck)", "0.9924", "0.9297", "16 / 162", "0.9886", "0.8707"],
    ["Focal, oversampling off  (selected)", "0.9906", "0.9374", "10 / 162", "0.9957", "0.8852"],
    ["Asymmetric focal, gamma 2.0, delta 0.60", "0.9916", "0.9201", "19 / 162", "0.9963", "0.8941"],
    ["Asymmetric focal, gamma 3.0, delta 0.75", "0.9872", "0.9224", "20 / 162", "0.9802", "0.8907"],
], top=1.55, col_w=[4.9, 1.4, 1.7, 1.5, 1.4, 1.7], font=11.5, height=2.1,
   highlight=(2,), hcolor=RGBColor(0xDF, 0xEA, 0xDF))

bullets(s, [
    "The gain came from switching off the oversampler, not from changing the loss. The sampler raised the "
    "positive rate the loss actually saw from 6.6 percent to 61 percent, inverting the imbalance it was "
    "meant to correct.",
    "Against the earlier baseline the DR gain is 0.0077 class-macro F1, which sits inside the noise "
    "(95% CI −0.006 to +0.022, p = 0.29). The claim that survives is the count: missed cases fell from "
    "16 to 10 out of 162.",
    "Scope. I ran two asymmetric-focal settings at the full 40 epochs against a matched control. "
    "The wider grid is designed but not run.",
], top=4.05, size=13, bottom=6.35)

takeaway(s, "Part 1 is closed. The strongest configuration is plain focal loss with the oversampler "
            "switched off.", GOOD, top=6.45)

# ═════════════════════════════════════════════ 2. the asymmetric focal question
s = slide()
title(s, "Part 1: did asymmetric focal loss help?",
      "Every metric, both losses. Oversampling is off in both arms, so the loss is the only difference.")

table(s, [
    ["Diabetic retinopathy (162 positive)", "Asym. focal", "Focal"],
    ["AUC", "0.9916", "0.9906"],
    ["Precision, diseased class", "0.8218", "0.8352"],
    ["Recall, diseased class", "0.8827", "0.9383"],
    ["F1, diseased class only", "0.8512", "0.8837"],
    ["F1, healthy class", "0.9890", "0.9912"],
    ["F1, class-macro (headline)", "0.9201", "0.9374"],
    ["Best F1-macro at any cutoff", "0.9292", "0.9374"],
    ["Cases missed", "19 / 162", "10 / 162"],
], top=1.45, left=0.65, width=6.0, col_w=[3.4, 1.3, 1.3],
   font=10, height=2.9, highlight=(6,), hcolor=RGBColor(0xF7, 0xE2, 0xDD))
table(s, [
    ["Macular edema (61 positive)", "Asym. focal", "Focal"],
    ["AUC", "0.9963", "0.9957"],
    ["Precision, diseased class", "0.8000", "0.8182"],
    ["Recall, diseased class", "0.7869", "0.7377"],
    ["F1, diseased class only", "0.7934", "0.7759"],
    ["F1, healthy class", "0.9947", "0.9945"],
    ["F1, class-macro (headline)", "0.8941", "0.8852"],
    ["Best F1-macro at any cutoff", "0.9022", "0.8995"],
    ["Cases missed", "13 / 61", "16 / 61"],
], top=1.45, left=7.0, width=5.75, col_w=[3.2, 1.3, 1.3],
   font=10, height=2.9, highlight=(6,), hcolor=RGBColor(0xDF, 0xEA, 0xDF))

label(s, 0.65, 4.48, 12.1,
      "Reading the two F1 rows. Class-macro F1 is the average of the diseased and healthy class F1 scores. "
      "For DR asymmetric focal that is (0.8512 + 0.9890) / 2 = 0.9201.",
      12, ACCENT, True, PP_ALIGN.LEFT)
label(s, 0.65, 4.84, 12.1,
      "Precision and recall are diseased-class only, so they sit well below the class-macro headline. "
      "Different denominators, not a contradiction.",
      11.5, MUTED, False, PP_ALIGN.LEFT)

bullets(s, [
    "Diabetic retinopathy. Asymmetric focal is worse by 0.0175 class-macro F1, 95% CI [−0.032, −0.004], "
    "p = 0.013. Its ranking is marginally better (AUC +0.0010), but even at the best cutoff available to it "
    "it reaches only 0.9292 against focal's 0.9374. This is not a threshold artifact.",
    "Macular edema. Asymmetric focal leads by 0.0090, but the interval [−0.012, +0.033] crosses zero "
    "(p = 0.448). I report this as no measurable difference, not as a win.",
], top=5.30, size=12.5, bottom=6.45)

takeaway(s, "Asymmetric focal loss did not beat focal loss on BRSET. It is significantly worse on "
            "diabetic retinopathy and shows no measurable difference on macular edema.", WARN, top=6.5)

# ═════════════════════════════════════════════ 3. exact configuration
s = slide()
title(s, "Part 1: exact configuration of every run",
      "Read from the run logs. Significance by paired bootstrap, 4,000 resamples, same test set.")

label(s, 0.65, 1.38, 6.0, "Identical across all four runs", 12.5, ACCENT, True, PP_ALIGN.LEFT)
table(s, [
    ["Setting", "Value"],
    ["Backbone", "ConvNeXt V2 Large, fcmae_ft_in22k_in1k_384"],
    ["Resolution", "512 px crop from a 560 px resize"],
    ["Optimizer", "AdamW, lr 5e-5, weight decay 0.1"],
    ["Schedule", "3-epoch linear warmup at 0.1x, then cosine to zero"],
    ["Epochs, effective batch", "40,  64"],
    ["Regularization", "drop-path 0.3, mixup alpha 0.2, label smoothing 0.1"],
    ["Augmentation", "random crop, horizontal flip, 15 degree rotation,\ncolour jitter 0.2 / 0.2 / 0.1"],
    ["Weight averaging", "EMA, decay 0.999"],
    ["Inference", "4-way flip TTA. Cutoff chosen on validation by\n200-resample bootstrap"],
    ["Precision, seed", "fp16 autocast,  seed 0"],
], top=1.72, left=0.65, width=6.0, col_w=[1.9, 4.1], font=9.5, height=3.05)

label(s, 7.0, 1.38, 5.75, "The only things that differ", 12.5, ACCENT, True, PP_ALIGN.LEFT)
table(s, [
    ["Run", "Loss", "Parameters", "Over-\nsample", "Positive rate\nseen by loss"],
    ["Baseline (earlier deck)", "Focal", "gamma 2.0", "on", "DR 61.0%\nME 38.5%"],
    ["Focal control", "Focal", "gamma 2.0", "off", "DR 6.6%\nME 2.4%"],
    ["Asymmetric focal A", "AFL", "gamma 2.0\ndelta 0.60", "off", "DR 6.6%\nME 2.4%"],
    ["Asymmetric focal B", "AFL", "gamma 3.0\ndelta 0.75", "off", "DR 6.6%\nME 2.4%"],
], top=1.72, left=7.0, width=5.75, col_w=[1.6, 0.8, 1.15, 0.8, 1.4], font=9.5, height=2.0)

box(s, 7.0, 3.88, 5.75, 0.40,
    "Focal:   L = − y (1−p)^g log p − (1−y) p^g log(1−p)", PALE, INK, 10.5, False, line=RULE)
box(s, 7.0, 4.34, 5.75, 0.40,
    "AFL:     L = − d y log p − (1−d)(1−y) p^g log(1−p)", PALE, INK, 10.5, False, line=RULE)
label(s, 7.0, 4.80, 5.75,
      "AFL differs in one way. Positives receive plain cross-entropy, with no (1−p)^g term.",
      10.5, MUTED, False, PP_ALIGN.LEFT, h=0.44)

label(s, 0.65, 5.30, 12.1, "What ran, and what did not", 12.5, ACCENT, True, PP_ALIGN.LEFT)
for head, col, txt, y in [
        ("Complete.", GOOD, "Three runs at the full 40 epochs, four Stage-0 controls at 20 epochs, "
                            "and the earlier baseline.", 5.66),
        ("Not run.", WARN, "Two of six planned grid points. Missing: gamma 1, 3 and 4 at delta 0.60, "
                           "and gamma 2 at delta 0.75.", 6.00),
        ("Cause.", MUTED, "Two crashed on a filesystem read error, since fixed. The rest were cancelled "
                          "when the design moved to 40 epochs.", 6.34)]:
    label(s, 0.65, y, 1.05, head, 11.5, col, True, PP_ALIGN.LEFT)
    label(s, 1.75, y, 10.95, txt, 11.5, INK, False, PP_ALIGN.LEFT)

takeaway(s, "Two settings at full length against a matched control. A controlled comparison, "
            "not a completed sweep.", AMBER, top=6.72)

# ═════════════════════════════════════════════ 4. Part 2, the problem
s = slide()
title(s, "Part 2: the model does not survive the change of camera",
      "Same disease and the same grading protocol. Only the camera and the population change.")

label(s, 0.65, 1.38, 6.5, "The two datasets", 12.5, ACCENT, True, PP_ALIGN.LEFT)
table(s, [
    ["", "BRSET, trained on", "mBRSET, tested on"],
    ["Camera", "Tabletop fundus camera", "Handheld phone camera\n(Phelcom Eyer)"],
    ["Setting", "Hospital eye clinic", "Community screening campaign"],
    ["Images used", "16,258  (test 2,435)", "4,859  (test 732)"],
    ["Patients with DR", "6.6%", "23.3%"],
    ["Images with an artifact", "not recorded", "83%  (4,272 of 5,164)"],
], top=1.72, left=0.65, width=6.5, col_w=[1.75, 2.35, 2.4], font=10, height=2.25)

label(s, 7.55, 1.38, 5.2, "One model, one cutoff of 0.61, two test sets", 12.5, ACCENT, True, PP_ALIGN.LEFT)
table(s, [
    ["Diabetic retinopathy", "on BRSET", "on mBRSET"],
    ["AUC", "0.988", "0.909"],
    ["Precision", "0.861", "0.964"],
    ["Recall", "0.877", "0.503"],
    ["F1", "0.869", "0.661"],
], top=1.72, left=7.55, width=5.2, col_w=[2.0, 1.6, 1.6], font=10,
   height=1.9, highlight=(3,), hcolor=RGBColor(0xF7, 0xE2, 0xDD))

label(s, 7.55, 3.70, 5.2, "Precision rose. Recall halved.", 12.5, WARN, True, PP_ALIGN.LEFT)

X0, XW, YB = 1.15, 11.0, 4.60
def sx(v):
    return X0 + XW * v

box(s, X0, YB, XW, 0.46, "", RGBColor(0xE7, 0xEA, 0xEF), INK, 10, False)
box(s, sx(0.19), YB, sx(0.61) - sx(0.19), 0.46, "79 diseased patients scored in this band",
    RGBColor(0xE0, 0xAE, 0x9C), INK, 11.5, True)
for v, txt, col in [(0.19, "0.19\nright cutoff for mBRSET", GOOD),
                    (0.61, "0.61\ncutoff learned on BRSET", WARN)]:
    box(s, sx(v) - 0.012, YB - 0.14, 0.024, 0.74, "", col, WHITE, 8, False)
    label(s, sx(v) - 1.35, YB + 0.52, 2.7, txt, 10.5, col, True, h=0.56)
label(s, X0, YB - 0.40, 1.4, "score 0", 10, MUTED, False, PP_ALIGN.LEFT)
label(s, X0 + XW - 1.4, YB - 0.40, 1.4, "score 1", 10, MUTED, False, PP_ALIGN.RIGHT)

label(s, 0.65, 5.86, 12.1,
      "BRSET is a clinic where 1 patient in 15 has diabetic retinopathy, so the model learned a strict "
      "cutoff of 0.61. mBRSET is a screening campaign where 1 in 4 does. Lowering the cutoff to 0.19, "
      "with no retraining, recovers 44 of those 79 patients.",
      12, INK, False, PP_ALIGN.LEFT, h=0.60)

takeaway(s, "The model still ranks patients well. It is the decision cutoff that moved. Measured with the "
            "earlier two-model ensemble, not yet re-run with the new Part 1 model.", ACCENT, top=6.55)

# ═════════════════════════════════════════════ 5. the decomposition figure
s = slide()
title(s, "Part 2: what the remaining gap is made of",
      "Diabetic retinopathy, F1 on mBRSET. The transferred model scores 0.661. A model trained directly "
      "on mBRSET scores 0.815. Every value measured on my own data.")

BX0, BW, BY, BH = 1.15, 11.0, 3.02, 0.62
LO, HI = 0.661, 0.815
def bx(v):
    return BX0 + BW * (v - LO) / (HI - LO)

SEGS = [(0.661, 0.717, C_DONE,  WHITE, "+0.056", "Recovered by retuning\nthe cutoff"),
        (0.717, 0.763, C_REACH, INK,   "+0.046", "Still reachable by a cutoff,\nbut not findable without\ntarget labels"),
        (0.763, 0.815, C_FEAT,  WHITE, "+0.052", "No cutoff reaches this.\nRequires better features")]

for a, b, fill, fg, delta, cap in SEGS:
    box(s, bx(a), BY, bx(b) - bx(a), BH, delta, fill, fg, 15, True, line=WHITE)
    label(s, bx(a), BY + BH + 0.12, bx(b) - bx(a), cap, 11, INK, False, h=0.70)

for v, cap in [(0.661, "0.661\ntransferred as is"), (0.717, "0.717\nafter retuning"),
               (0.763, "0.763\nbest at any cutoff"), (0.815, "0.815\ntrained on mBRSET")]:
    label(s, bx(v) - 1.0, 1.72, 2.0, cap, 11.5, ACCENT, True, h=0.60)
    box(s, bx(v) - 0.006, 2.36, 0.012, BY - 2.36, "", RULE, INK, 8, False)

for a, b, col, cap in [(0.661, 0.763, ACCENT, "Decision problem:  0.102 of the gap,  66 percent"),
                       (0.763, 0.815, C_FEAT, "Feature problem:  0.052,  34 percent")]:
    box(s, bx(a), 4.52, bx(b) - bx(a), 0.05, "", col, WHITE, 8, False)
    label(s, bx(a), 4.60, bx(b) - bx(a), cap, 11.5, col, True)

label(s, 0.65, 5.18, 5.9, "Why 0.763 is a hard ceiling", 12.5, ACCENT, True, PP_ALIGN.LEFT)
label(s, 0.65, 5.50, 5.9,
      "I swept every cutoff from 0.01 to 0.99. The best F1 attainable is 0.763. Calibration, temperature "
      "scaling and label-shift correction only rescale scores, so each picks a point on this same curve. "
      "None of them can pass 0.763, which rules out that whole family.",
      11.5, INK, False, PP_ALIGN.LEFT, h=1.15)

label(s, 7.0, 5.18, 5.75, "Why this is not only a prevalence change", 12.5, ACCENT, True, PP_ALIGN.LEFT)
label(s, 7.0, 5.50, 5.75,
      "ROC is invariant to how many patients are diseased (Fawcett 2006). If only the disease rate had "
      "changed, AUC would have held steady. It fell from 0.988 to 0.909, so the images themselves differ. "
      "The handheld camera changed what the model sees.",
      11.5, INK, False, PP_ALIGN.LEFT, h=1.15)

takeaway(s, "Two thirds of the gap is a decision problem and one third is a feature problem. Most "
            "cross-device papers report a single number and blame the model. This split is computable.",
         top=6.70)

# ═════════════════════════════════════════════ 6. literature
s = slide()
title(s, "Part 2: what the published literature establishes",
      "Peer-reviewed versions only, never the preprint. Citation counts from Semantic Scholar, "
      "7 August 2026. Paper names are live links.")

lit = [
    ["Evidence", "Finding", "Paper", "Venue", "Cites"],
    ["Supports", "Appearance augmentation beat every adaptation\nmethod tested for cross-scanner shift",
     "Tellez et al.", "Medical Image\nAnalysis 2019", "656"],
    ["Supports", "Fundus DR over six datasets, AUC 75.9 to 82.6.\nAugmentation alone contributes 3.4",
     "Che et al.\n(GDRNet)", "MICCAI 2023", "46"],
    ["Fails", "DANN, MDD and CycleGAN all collapse on real\ntabletop to portable fundus data",
     "Lin et al.", "MICCAI 2022", "5"],
    ["Fails", "No domain-generalization method reliably beats\nplain training under fair evaluation",
     "Gulrajani and\nLopez-Paz", "ICLR 2021", "1,504"],
    ["Fails", "Test-time adaptation loses up to 66 points under\nprior shift, the exact regime here",
     "Boudiaf et al.\n(LAME)", "CVPR 2022", "250"],
    ["Partial", "Cutoff can be fixed with no target labels, but the\nmethod assumes prevalence is unchanged",
     "Roschewitz et al.", "Nature\nComms 2023", "42"],
]
shp = table(s, lit, top=1.62, col_w=[0.9, 5.2, 2.0, 2.0, 0.7], font=10.5, height=3.5)
for row, (col, url) in enumerate([
        (GOOD,  "https://doi.org/10.1016/j.media.2019.101544"),
        (GOOD,  "https://doi.org/10.1007/978-3-031-43904-9_42"),
        (WARN,  "https://doi.org/10.1007/978-3-031-16434-7_57"),
        (WARN,  "https://openreview.net/forum?id=lQdXeXDoWtI"),
        (WARN,  "https://doi.org/10.1109/CVPR52688.2022.00816"),
        (AMBER, "https://doi.org/10.1038/s41467-023-42396-y")], start=1):
    cell_color(shp, row, 0, col)
    link_cell(shp, row, 2, url)

bullets(s, [
    "The three methods a reader would reach for first, adversarial adaptation, domain generalization and "
    "test-time adaptation, are all published failures in exactly this setting.",
    "Augmentation is the only intervention with peer-reviewed evidence of raising AUC under real fundus "
    "device shift, which is why it is the direction I propose.",
], top=5.35, size=12.5, bottom=6.5)

takeaway(s, "This review changed the plan. It removed three candidate methods and left one.", top=6.55)

# ═════════════════════════════════════════════ 7. the falsified prediction
s = slide()
title(s, "Part 2: I tested my own proposal before building it",
      "mBRSET gives about four images per patient. Theory says pooling them should raise AUC, and it "
      "reduces to a single measurable quantity, so the idea can be killed in an afternoon.")

label(s, 0.65, 1.36, 12.1,
      "What rho measures. Each patient has about four images. Rho asks whether those are four independent "
      "opinions, or one opinion repeated four times. It is the intraclass correlation "
      "(Shrout and Fleiss 1979 [12]).",
      12, INK, False, PP_ALIGN.LEFT, h=0.56)

RX0, RXW, RY = 2.70, 7.90, 2.40
box(s, RX0, RY, RXW, 0.26, "", RGBColor(0xE2, 0xE6, 0xEC), INK, 8, False)
box(s, RX0 + RXW * 0.363 - 0.014, RY - 0.10, 0.028, 0.46, "", ACCENT, WHITE, 8, False)
label(s, RX0 + RXW * 0.363 - 1.85, RY - 0.42, 3.7,
      "measured for diabetic retinopathy:  rho = 0.363", 11, ACCENT, True)
label(s, RX0 - 0.32, RY - 0.03, 0.3, "0", 10.5, MUTED, True)
label(s, RX0 + RXW + 0.04, RY - 0.03, 0.3, "1", 10.5, MUTED, True)
label(s, RX0, RY + 0.32, 3.7,
      "Images disagree freely. Pooling four\ncuts noise like four separate readings.",
      10.5, INK, False, PP_ALIGN.LEFT, h=0.48)
label(s, RX0 + RXW - 3.7, RY + 0.32, 3.7,
      "Images all say the same thing.\nPooling four adds nothing.",
      10.5, INK, False, PP_ALIGN.RIGHT, h=0.48)

box(s, 0.65, 3.32, 5.4, 0.44, "AUC after pooling K  =  Phi( z1 * sqrt( K / (1 + (K−1) rho) ) )",
    PALE, INK, 11.5, True, line=RULE)
label(s, 6.30, 3.34, 6.45,
      "Rho is the only unknown. Measure it once on data already in hand and the achievable gain follows, "
      "with no training run.",
      11.5, INK, False, PP_ALIGN.LEFT, h=0.44)

shp = table(s, [
    ["What I pooled", "rho", "AUC predicted", "AUC measured", "Outcome"],
    ["DR, four images per patient, mean", "0.363", "0.968", "0.844", "prediction failed"],
    ["DR, two images per eye, mean", "0.481", "0.940", "0.885", "prediction failed"],
    ["ME, four images per patient, max", "0.309", "not applicable", "0.969", "worked, plus 0.036 AUC"],
], top=3.90, left=0.65, width=12.1, col_w=[3.8, 1.0, 2.0, 1.9, 2.5], font=11, height=1.30,
   highlight=(3,), hcolor=RGBColor(0xDF, 0xEA, 0xDF))
cell_color(shp, 1, 4, WARN); cell_color(shp, 2, 4, WARN); cell_color(shp, 3, 4, GOOD)
label(s, 0.65, 5.24, 12.1,
      "Image-level AUC before any pooling: DR 0.909, ME 0.933.  Pooling operators: Ilse et al., ICML 2018 [13].",
      10.5, MUTED, False, PP_ALIGN.LEFT)

bullets(s, [
    "Why it failed. Rho was measured correctly. The hidden assumption is what broke: the formula treats "
    "every image in a group as carrying the same finding. Of the 187 patients with both eyes imaged, 22 "
    "of them, 11.8 percent, have one eye affected and one clear. Averaging a diseased eye with a healthy "
    "one lands between the two, which is worse than either.",
    "Macular edema is focal, so a maximum is the right operator rather than a mean. That works, lifting "
    "AUC from 0.933 to 0.969.",
], top=5.60, size=11.5, bottom=6.68)

takeaway(s, "Rho let me price the idea before building it. The prediction was wrong, and the reason is "
            "now understood.", WARN, top=6.72)

# ═════════════════════════════════════════════ 8. the opening
s = slide()
title(s, "Part 2: where the opening is, and what I propose")

label(s, 0.65, 1.42, 12.1, "The nearest existing work, and what each one leaves open",
      12.5, ACCENT, True, PP_ALIGN.LEFT)
shp = table(s, [
    ["Work", "What it does", "What it leaves open"],
    ["RetSyn, J. Biomedical Informatics 2025\n(same senior author as both datasets)",
     "Synthetic data for tabletop to portable",
     "Does not diagnose what breaks. No macular edema.\nNo treatment of the decision cutoff"],
    ["IEEE ISBI 2026", "Reports that BRSET to mBRSET drops",
     "Descriptive only. Not what the drop is made of,\nnor how to close it"],
], top=1.78, col_w=[3.9, 3.3, 4.9], font=10.5, height=1.4)
link_cell(shp, 1, 0, "https://doi.org/10.1016/j.jbi.2025.104938")

label(s, 0.65, 3.42, 12.1, "Three openings this review did not find addressed anywhere",
      12.5, ACCENT, True, PP_ALIGN.LEFT)
y = 3.82
for n, t in [("1", "Splitting a cross-device gap into prevalence, cutoff and ranking, and then proving the "
                   "cutoff family is exhausted at 0.763, rather than trying methods one at a time."),
             ("2", "Label-shift correction for multi-label sigmoid outputs. BBSE (ICML 2018) and MLLS "
                   "(ICML 2020, NeurIPS 2020) are all single-label softmax. DR and macular edema co-occur, "
                   "so correcting each label on its own is misspecified."),
             ("3", "The regime where both published families break their own assumptions at once. One "
                   "assumes prevalence is fixed, and here it moves from 6.6 to 23.3 percent. The other "
                   "assumes the images look the same, and here AUC fell from 0.988 to 0.909.")]:
    box(s, 0.65, y, 0.40, 0.40, n, ACCENT, WHITE, 12.5, True)
    label(s, 1.20, y - 0.04, 11.5, t, 12, INK, False, PP_ALIGN.LEFT, h=0.74)
    y += 0.84

takeaway(s, "Proposed next step: fundus degradation augmentation, the one intervention with published "
            "evidence here, with the gap decomposition as the scientific contribution.", GOOD, top=6.45)

# ═════════════════════════════════════════════ 9. references
s = slide()
title(s, "References",
      "Peer-reviewed versions. Citation counts from Semantic Scholar, 7 August 2026. All entries are live links.")
refs = [
 ("Tellez et al. Quantifying the effects of data augmentation and stain color normalization in convolutional neural networks for computational pathology. Medical Image Analysis 58:101544, 2019. 656 cites.",
  "https://doi.org/10.1016/j.media.2019.101544"),
 ("Che, Cheng, Jin, Chen. Towards Generalizable Diabetic Retinopathy Grading in Unseen Domains. MICCAI 2023, 430-440. 46 cites.",
  "https://doi.org/10.1007/978-3-031-43904-9_42"),
 ("Lin, Shi, Zhang, Shang, He, Ge. Camera Adaptation for Fundus-Image-Based CVD Risk Estimation. MICCAI 2022, 593-603. 5 cites.",
  "https://doi.org/10.1007/978-3-031-16434-7_57"),
 ("Gulrajani, Lopez-Paz. In Search of Lost Domain Generalization. ICLR 2021. 1,504 cites.",
  "https://openreview.net/forum?id=lQdXeXDoWtI"),
 ("Boudiaf, Mueller, Ben Ayed, Bertinetto. Parameter-free Online Test-time Adaptation. CVPR 2022, 8344-8353. 250 cites.",
  "https://doi.org/10.1109/CVPR52688.2022.00816"),
 ("Roschewitz, Khara, Yearsley et al. Automatic correction of performance drift under acquisition shift in medical image classification. Nature Communications 14:6608, 2023. 42 cites.",
  "https://doi.org/10.1038/s41467-023-42396-y"),
 ("Godau, Kalinowski, Christodoulou et al. Navigating prevalence shifts in image analysis algorithm deployment. Medical Image Analysis 102:103504, 2025.",
  "https://doi.org/10.1016/j.media.2025.103504"),
 ("Lipton, Wang, Smola. Detecting and Correcting for Label Shift with Black Box Predictors. ICML 2018. 679 cites.",
  "https://proceedings.mlr.press/v80/lipton18a.html"),
 ("Alexandari, Kundaje, Shrikumar. Maximum Likelihood with Bias-Corrected Calibration is Hard-To-Beat at Label Shift Adaptation. ICML 2020, 222-232.",
  "https://proceedings.mlr.press/v119/alexandari20a.html"),
 ("Garg, Wu, Balakrishnan, Lipton. A Unified View of Label Shift Estimation. NeurIPS 2020. 179 cites.",
  "https://proceedings.neurips.cc/paper/2020/hash/219e052492f4008818b8adb6366c7ed6-Abstract.html"),
 ("Fawcett. An introduction to ROC analysis. Pattern Recognition Letters 27(8):861-874, 2006. 21,726 cites.",
  "https://doi.org/10.1016/j.patrec.2005.10.010"),
 ("Shrout, Fleiss. Intraclass correlations: uses in assessing rater reliability. Psychological Bulletin 86(2):420-428, 1979.",
  "https://doi.org/10.1037/0033-2909.86.2.420"),
 ("Ilse, Tomczak, Welling. Attention-based Deep Multiple Instance Learning. ICML 2018. 2,691 cites.",
  "https://proceedings.mlr.press/v80/ilse18a.html"),
 ("Shuai, Wu, Tang, Restrepo, Morley, Nakayama. Enhancing AI-based diabetic retinopathy screening in LMICs with synthetic data (RetSyn). Journal of Biomedical Informatics 172:104938, 2025.",
  "https://doi.org/10.1016/j.jbi.2025.104938"),
 ("Nakayama, Restrepo et al. BRSET: A Brazilian Multilabel Ophthalmological Dataset. PLOS Digital Health 3(7):e0000454, 2024.",
  "https://doi.org/10.1371/journal.pdig.0000454"),
 ("Wu, Restrepo, Nakayama et al. A portable retina fundus photos dataset for clinical, demographic, and diabetic retinopathy prediction (mBRSET). Scientific Data 12:340, 2025.",
  "https://doi.org/10.1038/s41597-025-04672-y"),
]
b = s.shapes.add_textbox(Inches(0.65), Inches(1.50), Inches(12.1), Inches(5.65))
tf = b.text_frame; tf.word_wrap = True
for i, (txt, url) in enumerate(refs):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    r = p.add_run(); r.text = f"[{i+1}]  {txt}  "
    r.font.size = Pt(9); r.font.color.rgb = INK; r.font.name = FONT
    lr = p.add_run(); lr.text = url
    lr.font.size = Pt(9); lr.font.name = FONT
    lr.hyperlink.address = url
    p.space_after = Pt(2)

prs.save(OUT)
print(f"wrote {OUT}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
